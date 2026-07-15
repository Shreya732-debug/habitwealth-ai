# backend/rag_engine.py
"""
RAG (Retrieval-Augmented Generation) engine for FinanceGPT.

Supported file types:
  - PDF                   (.pdf)
  - Word documents        (.docx, .doc)
  - PowerPoint            (.pptx, .ppt)
  - Excel spreadsheets    (.xlsx, .xls)
  - CSV files             (.csv)
  - Images / screenshots  (.png, .jpg, .jpeg, .webp, .bmp, .tiff)
  - Plain text            (.txt)

Pipeline:
  1. File upload → format detection → text extraction
  2. Text chunking (500 chars, 50 overlap)
  3. Embedding generation (sentence-transformers, local, free)
  4. Vector storage (ChromaDB, local, free)
  5. Query → similarity search → relevant chunks → injected into Gemini prompt
"""

import os
import re
import csv
import hashlib
import tempfile
from typing import List, Dict, Optional

import pdfplumber
import chromadb
from sentence_transformers import SentenceTransformer

# ── Constants ─────────────────────────────────────────────────
CHROMA_PATH     = "./chroma_db"
COLLECTION_NAME = "finance_docs"
CHUNK_SIZE      = 500
CHUNK_OVERLAP   = 50
TOP_K_RESULTS   = 3
MAX_FILE_MB     = 15

SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf",
    ".docx", ".doc",
    ".pptx", ".ppt",
    # Spreadsheets
    ".xlsx", ".xls",
    ".csv",
    # Images / screenshots
    ".png", ".jpg", ".jpeg",
    ".webp", ".bmp", ".tiff", ".tif",
    # Plain text
    ".txt",
}

# ── Lazy-loaded singletons ────────────────────────────────────
_embedding_model = None
_collection      = None


def _get_model() -> SentenceTransformer:
    global _embedding_model
    if _embedding_model is None:
        print("[RAG] Loading embedding model (first time only)...")
        _embedding_model = SentenceTransformer("all-MiniLM-L6-v2")
        print("[RAG] Model ready.")
    return _embedding_model


def _get_collection():
    global _collection
    if _collection is None:
        client = chromadb.PersistentClient(path=CHROMA_PATH)
        _collection = client.get_or_create_collection(
            name=COLLECTION_NAME,
            metadata={"hnsw:space": "cosine"}
        )
    return _collection


# ══════════════════════════════════════════════════════════════
#  TEXT EXTRACTORS — one per file type
# ══════════════════════════════════════════════════════════════

def _extract_pdf(path: str) -> List[Dict]:
    """Extract text from PDF, page by page."""
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append({"page_num": i + 1, "text": text.strip()})

            # Also extract tables from PDF pages
            tables = page.extract_tables()
            for table in tables:
                if table:
                    rows = []
                    for row in table:
                        cleaned = [str(cell).strip() for cell in row if cell]
                        if cleaned:
                            rows.append(" | ".join(cleaned))
                    if rows:
                        table_text = "\n".join(rows)
                        pages.append({
                            "page_num": i + 1,
                            "text": f"[Table on page {i+1}]\n{table_text}"
                        })
    return pages


def _extract_docx(path: str) -> List[Dict]:
    """Extract text from Word documents (.docx)."""
    from docx import Document
    doc = Document(path)
    pages = []
    current_text = []
    page_num = 1

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            current_text.append(text)
        # Simulate page breaks every 40 paragraphs
        if len(current_text) >= 40:
            pages.append({"page_num": page_num, "text": "\n".join(current_text)})
            current_text = []
            page_num += 1

    # Extract tables from Word doc
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            current_text.append("[Table]\n" + "\n".join(rows))

    if current_text:
        pages.append({"page_num": page_num, "text": "\n".join(current_text)})

    return pages


def _extract_pptx(path: str) -> List[Dict]:
    """Extract text from PowerPoint presentations (.pptx)."""
    from pptx import Presentation
    prs = Presentation(path)
    pages = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []

        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

            # Extract table data from slides
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        rows.append(" | ".join(cells))
                if rows:
                    texts.append("[Table]\n" + "\n".join(rows))

        if texts:
            pages.append({
                "page_num": slide_num,
                "text": f"[Slide {slide_num}]\n" + "\n".join(texts)
            })

    return pages


def _extract_xlsx(path: str) -> List[Dict]:
    """Extract text from Excel spreadsheets (.xlsx, .xls)."""
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    pages = []

    for sheet_num, sheet_name in enumerate(wb.sheetnames, 1):
        ws = wb[sheet_name]
        rows = []

        # Get header row first
        headers = []
        for row in ws.iter_rows(max_row=1, values_only=True):
            headers = [str(cell).strip() if cell is not None else "" for cell in row]

        # Process all rows
        for row_idx, row in enumerate(ws.iter_rows(values_only=True), 1):
            cells = []
            for col_idx, cell in enumerate(row):
                if cell is not None and str(cell).strip():
                    # Include header context if available
                    header = headers[col_idx] if col_idx < len(headers) else f"Col{col_idx}"
                    cells.append(f"{header}: {str(cell).strip()}")
            if cells:
                rows.append(", ".join(cells))

        if rows:
            pages.append({
                "page_num": sheet_num,
                "text": f"[Sheet: {sheet_name}]\n" + "\n".join(rows[:500])  # cap at 500 rows
            })

    return pages


def _extract_csv(path: str) -> List[Dict]:
    """Extract text from CSV files."""
    pages = []
    rows = []

    # Try multiple encodings
    for encoding in ["utf-8", "latin-1", "cp1252"]:
        try:
            with open(path, "r", encoding=encoding) as f:
                reader = csv.reader(f)
                headers = []
                for row_idx, row in enumerate(reader):
                    if row_idx == 0:
                        headers = row
                        rows.append("Headers: " + " | ".join(headers))
                    else:
                        if any(cell.strip() for cell in row):
                            # Format as "header: value" pairs
                            pairs = []
                            for i, cell in enumerate(row):
                                if cell.strip():
                                    h = headers[i] if i < len(headers) else f"Col{i}"
                                    pairs.append(f"{h}: {cell.strip()}")
                            rows.append(", ".join(pairs))
            break
        except UnicodeDecodeError:
            continue

    if rows:
        # Split into pages of 100 rows each
        page_size = 100
        for i in range(0, len(rows), page_size):
            chunk = rows[i:i + page_size]
            pages.append({
                "page_num": (i // page_size) + 1,
                "text": "\n".join(chunk)
            })

    return pages


def _extract_image(path: str) -> List[Dict]:
    """
    Extract text from images and screenshots using OCR (Tesseract).
    Works for: screenshots of bank statements, photos of receipts,
               scanned documents, etc.
    """
    try:
        import pytesseract
        from PIL import Image

        # Windows — set Tesseract path if not in PATH
        if os.name == 'nt':
            tesseract_paths = [
                r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            ]
            for tp in tesseract_paths:
                if os.path.exists(tp):
                    pytesseract.pytesseract.tesseract_cmd = tp
                    break

        img = Image.open(path)

        # Enhance image for better OCR accuracy
        # Convert to RGB if needed (handles PNG with transparency)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        # Run OCR
        text = pytesseract.image_to_string(img, lang="eng")
        text = text.strip()

        if not text:
            return [{"page_num": 1, "text": "[Image contained no readable text]"}]

        return [{"page_num": 1, "text": text}]

    except ImportError:
        raise ValueError(
            "OCR not available. Install pytesseract and Tesseract:\n"
            "pip install pytesseract pillow\n"
            "Then install Tesseract from: https://github.com/UB-Mannheim/tesseract/wiki"
        )
    except Exception as e:
        raise ValueError(f"Image OCR failed: {str(e)}")


def _extract_txt(path: str) -> List[Dict]:
    """Extract text from plain text files."""
    for encoding in ["utf-8", "latin-1", "cp1252"]:
        try:
            with open(path, "r", encoding=encoding) as f:
                content = f.read()
            if content.strip():
                # Split into pages of ~2000 chars
                pages = []
                for i in range(0, len(content), 2000):
                    pages.append({
                        "page_num": (i // 2000) + 1,
                        "text": content[i:i + 2000].strip()
                    })
                return pages
        except UnicodeDecodeError:
            continue
    return []


# ── Dispatcher ────────────────────────────────────────────────

def extract_text(file_path: str, filename: str) -> List[Dict]:
    """
    Route to the correct extractor based on file extension.
    Returns list of {page_num, text} dicts.
    """
    ext = os.path.splitext(filename)[1].lower()

    extractor_map = {
        ".pdf":  _extract_pdf,
        ".docx": _extract_docx,
        ".doc":  _extract_docx,
        ".pptx": _extract_pptx,
        ".ppt":  _extract_pptx,
        ".xlsx": _extract_xlsx,
        ".xls":  _extract_xlsx,
        ".csv":  _extract_csv,
        ".txt":  _extract_txt,
        ".png":  _extract_image,
        ".jpg":  _extract_image,
        ".jpeg": _extract_image,
        ".webp": _extract_image,
        ".bmp":  _extract_image,
        ".tiff": _extract_image,
        ".tif":  _extract_image,
    }

    extractor = extractor_map.get(ext)
    if not extractor:
        raise ValueError(
            f"Unsupported file type: {ext}\n"
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    pages = extractor(file_path)

    if not pages:
        raise ValueError(
            f"No readable text found in '{filename}'. "
            f"If this is a scanned document, try converting it to a searchable PDF first."
        )

    # Clean extracted text
    for page in pages:
        page["text"] = re.sub(r'\s+', ' ', page["text"]).strip()

    return pages


# ══════════════════════════════════════════════════════════════
#  CHUNKING
# ══════════════════════════════════════════════════════════════

def chunk_pages(pages: List[Dict], source_name: str) -> List[Dict]:
    """Split pages into overlapping chunks for embedding."""
    chunks = []

    for page in pages:
        text     = page["text"]
        page_num = page["page_num"]
        start    = 0
        idx      = 0

        while start < len(text):
            end = start + CHUNK_SIZE

            # Try to end at a sentence boundary
            if end < len(text):
                for punct in [". ", "! ", "? ", "\n", ", "]:
                    boundary = text.rfind(punct, start, end)
                    if boundary > start + CHUNK_SIZE // 2:
                        end = boundary + 1
                        break

            chunk = text[start:end].strip()
            if len(chunk) > 20:  # skip tiny fragments
                chunks.append({
                    "text":     chunk,
                    "source":   source_name,
                    "page_num": page_num,
                    "chunk_id": hashlib.md5(
                        f"{source_name}_{page_num}_{idx}".encode()
                    ).hexdigest()
                })
                idx += 1

            start = end - CHUNK_OVERLAP
            if start >= len(text):
                break

    return chunks


# ══════════════════════════════════════════════════════════════
#  INDEX / STORE
# ══════════════════════════════════════════════════════════════

def index_document(
    file_path: str,
    user_id: str,
    original_filename: str
) -> Dict:
    """
    Full pipeline: file → extract → chunk → embed → store.
    Works for all supported file types.
    """
    # Extract text (routes to correct extractor)
    pages = extract_text(file_path, original_filename)

    # Chunk
    chunks = chunk_pages(pages, original_filename)
    if not chunks:
        raise ValueError("Could not extract any text chunks from this file.")

    # Embed all chunks at once (batch is faster)
    model      = _get_model()
    texts      = [c["text"] for c in chunks]
    embeddings = model.encode(texts, show_progress_bar=False).tolist()

    # Store in ChromaDB
    collection = _get_collection()
    ids        = [f"{user_id}_{c['chunk_id']}" for c in chunks]
    metadatas  = [
        {
            "user_id":   user_id,
            "source":    c["source"],
            "page_num":  str(c["page_num"]),
        }
        for c in chunks
    ]

    collection.upsert(
        ids=ids,
        embeddings=embeddings,
        documents=texts,
        metadatas=metadatas
    )

    ext = os.path.splitext(original_filename)[1].upper().lstrip(".")

    return {
        "filename":  original_filename,
        "file_type": ext,
        "pages":     len(pages),
        "chunks":    len(chunks),
        "indexed":   True,
        "message":   f"Indexed {len(chunks)} chunks from {len(pages)} page(s)."
    }


# ══════════════════════════════════════════════════════════════
#  RETRIEVE
# ══════════════════════════════════════════════════════════════

def retrieve_relevant_chunks(
    query: str,
    user_id: str,
    top_k: int = TOP_K_RESULTS
) -> List[Dict]:
    """Find the most relevant chunks for a query."""
    collection = _get_collection()

    # Check if user has any documents
    existing = collection.get(where={"user_id": user_id}, limit=1)
    if not existing["ids"]:
        return []

    model           = _get_model()
    query_embedding = model.encode([query], show_progress_bar=False).tolist()[0]

    results = collection.query(
        query_embeddings=[query_embedding],
        n_results=min(top_k, collection.count()),
        where={"user_id": user_id},
        include=["documents", "metadatas", "distances"]
    )

    if not results["ids"][0]:
        return []

    chunks = []
    for i in range(len(results["ids"][0])):
        distance   = results["distances"][0][i]
        similarity = 1 - (distance / 2)

        if similarity > 0.3:
            chunks.append({
                "text":       results["documents"][0][i],
                "source":     results["metadatas"][0][i]["source"],
                "page_num":   results["metadatas"][0][i]["page_num"],
                "similarity": round(similarity, 3)
            })

    return chunks


# ══════════════════════════════════════════════════════════════
#  LIST / DELETE
# ══════════════════════════════════════════════════════════════

def list_user_documents(user_id: str) -> List[Dict]:
    """List all documents with their file types."""
    collection = _get_collection()
    results    = collection.get(where={"user_id": user_id})

    if not results["metadatas"]:
        return []

    seen    = {}
    for meta in results["metadatas"]:
        src = meta["source"]
        if src not in seen:
            ext = os.path.splitext(src)[1].upper().lstrip(".") or "FILE"
            seen[src] = {"filename": src, "type": ext}

    return list(seen.values())


def delete_document(user_id: str, filename: str) -> Dict:
    """Delete all chunks for a document."""
    collection = _get_collection()
    results    = collection.get(where={"user_id": user_id, "source": filename})

    if not results["ids"]:
        return {"deleted": False, "message": "Document not found."}

    collection.delete(ids=results["ids"])
    return {
        "deleted":        True,
        "filename":       filename,
        "chunks_removed": len(results["ids"])
    }